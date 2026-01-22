# 1.2.42: PPT/PPTX 文件加载器 - 基于 Chatmax GeneralPPTXLoader 实现
# 支持文本提取、表格识别、图片OCR（可选）

import os
from typing import List, Optional
from io import BytesIO
from langchain_core.documents import Document

# 核心依赖
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 可选依赖：图片OCR处理
try:
    from PIL import Image, ImageEnhance
    import numpy as np
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import cv2
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# 1.2.42: 环境检测
ENV = os.getenv("ENV", "production")

def log_debug(message: str):
    """仅在开发环境输出日志"""
    if ENV == "development":
        print(message)


class GeneralPPTXLoader:
    """
    1.2.42: PPT/PPTX 文件加载器

    基于 Chatmax 的 GeneralPPTXLoader 实现，支持：
    - 文本提取：从所有形状中提取文本内容
    - 表格处理：提取表格数据并转换为文本
    - 图片OCR：提取图片中的文字（需要安装 pytesseract）
    - 分组形状：递归处理嵌套的形状组

    使用方法：
        loader = GeneralPPTXLoader("path/to/file.pptx")
        documents = loader.load()
    """

    def __init__(self, file_path: str, enable_ocr: bool = False):
        """
        初始化 PPTX 加载器

        Args:
            file_path: PPT 文件路径
            enable_ocr: 是否启用图片 OCR（需要安装 pytesseract 和 tesseract-ocr）
        """
        self.file_path = file_path
        self.enable_ocr = enable_ocr and HAS_PIL and HAS_CV2 and HAS_TESSERACT

        # 加载 PPT 文件
        try:
            self.prs = Presentation(file_path)
            log_debug(f"✅ 成功加载 PPTX 文件: {file_path}")
        except Exception as e:
            raise ValueError(f"无法加载 PPTX 文件 {file_path}: {str(e)}")

        # OCR 依赖检查
        if enable_ocr and not self.enable_ocr:
            missing = []
            if not HAS_PIL:
                missing.append("pillow")
            if not HAS_CV2:
                missing.append("opencv-python")
            if not HAS_TESSERACT:
                missing.append("pytesseract")
            log_debug(f"⚠️ OCR 功能已禁用，缺少依赖: {', '.join(missing)}")

    def get_slide_count(self) -> int:
        """获取幻灯片总数"""
        return len(self.prs.slides)

    def load(self) -> List[Document]:
        """
        加载 PPTX 文件并返回 Document 列表

        每张幻灯片生成一个 Document，包含：
        - page_content: 幻灯片中的所有文本内容
        - metadata: 文件来源、页码等元数据

        Returns:
            List[Document]: 文档列表
        """
        documents: List[Document] = []
        slide_count = self.get_slide_count()

        log_debug(f"📊 开始解析 PPTX，共 {slide_count} 张幻灯片")

        for slide_idx, slide in enumerate(self.prs.slides):
            # 解析幻灯片内容
            slide_content = self._parse_slide(slide, slide_idx)

            # 合并所有文本内容
            all_texts = []

            # 添加文本内容
            for text_item in slide_content.get('texts', []):
                content = text_item.get('content', '').strip()
                if content:
                    all_texts.append(content)

            # 添加表格内容（转换为文本）
            for table_item in slide_content.get('tables', []):
                table_text = self._table_to_text(table_item.get('content', []))
                if table_text:
                    all_texts.append(table_text)

            # 添加图片 OCR 内容
            if self.enable_ocr:
                for picture_item in slide_content.get('pictures', []):
                    ocr_text = picture_item.get('ocr_text', '').strip()
                    if ocr_text:
                        all_texts.append(f"[图片文字] {ocr_text}")

            # 合并为页面内容
            page_content = '\n\n'.join(all_texts)

            # 创建 Document
            doc = Document(
                page_content=page_content,
                metadata={
                    "source": self.file_path,
                    "file_path": self.file_path,
                    "file_type": "pptx",
                    "page_number": slide_idx + 1,
                    "total_pages": slide_count,
                    "has_tables": len(slide_content.get('tables', [])) > 0,
                    "has_pictures": len(slide_content.get('pictures', [])) > 0
                }
            )
            documents.append(doc)

            log_debug(f"  幻灯片 {slide_idx + 1}/{slide_count}: "
                     f"文本={len(slide_content.get('texts', []))}个, "
                     f"表格={len(slide_content.get('tables', []))}个, "
                     f"图片={len(slide_content.get('pictures', []))}个")

        log_debug(f"✅ PPTX 解析完成，生成 {len(documents)} 个文档")
        return documents

    def _parse_slide(self, slide, slide_idx: int) -> dict:
        """
        解析单张幻灯片

        Args:
            slide: 幻灯片对象
            slide_idx: 幻灯片索引

        Returns:
            dict: 包含 texts, tables, pictures 的字典
        """
        slide_content = {
            'texts': [],
            'tables': [],
            'pictures': []
        }

        for shape_idx, shape in enumerate(slide.shapes):
            self._parse_shape(shape, slide_idx, shape_idx, slide_content)

        return slide_content

    def _parse_shape(self, shape, slide_idx: int, shape_idx, slide_content: dict):
        """
        解析形状内容

        根据形状类型调用不同的处理方法：
        - GROUP: 递归处理分组
        - PICTURE: 提取图片并进行 OCR
        - TABLE: 提取表格数据
        - 其他有文本框的形状: 提取文本

        Args:
            shape: 形状对象
            slide_idx: 幻灯片索引
            shape_idx: 形状索引
            slide_content: 存储解析结果的字典
        """
        try:
            # 分组形状 - 递归处理
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._parse_group(shape, slide_idx, shape_idx, slide_content)

            # 图片 - 提取并可选 OCR
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_data = self._parse_picture(shape, slide_idx, shape_idx)
                if picture_data:
                    slide_content['pictures'].append(picture_data)

            # 表格 - 提取数据
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = self._parse_table(shape)
                if table_data:
                    slide_content['tables'].append(table_data)

            # 其他形状 - 尝试提取文本
            elif shape.has_text_frame:
                text_data = self._parse_text_frame(shape)
                if text_data and text_data.get('content'):
                    slide_content['texts'].append(text_data)

            # 自动形状和自由形状也可能有文本
            elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text_data = self._parse_text_frame(shape)
                    if text_data and text_data.get('content'):
                        slide_content['texts'].append(text_data)

        except Exception as e:
            log_debug(f"⚠️ 解析形状出错 (幻灯片 {slide_idx}, 形状 {shape_idx}): {str(e)}")

    def _parse_text_frame(self, shape) -> dict:
        """
        提取形状中的文本

        Args:
            shape: 形状对象

        Returns:
            dict: {"content": 文本内容}
        """
        try:
            text = shape.text.strip() if shape.text else ""
            return {"content": text}
        except Exception:
            return {"content": ""}

    def _parse_picture(self, shape, slide_idx: int, shape_idx) -> Optional[dict]:
        """
        处理图片形状

        如果启用 OCR，会对图片进行预处理和文字识别

        Args:
            shape: 图片形状对象
            slide_idx: 幻灯片索引
            shape_idx: 形状索引

        Returns:
            dict: {"ocr_text": OCR识别的文字} 或 None
        """
        if not self.enable_ocr:
            return {"ocr_text": ""}

        try:
            image = shape.image
            image_stream = BytesIO(image.blob)

            try:
                image_data = Image.open(image_stream)
            except OSError as e:
                # WMF 格式处理 - 跳过（需要 ImageMagick）
                if 'cannot find loader for this WMF file' in str(e):
                    log_debug(f"⚠️ 跳过 WMF 格式图片 (幻灯片 {slide_idx}, 形状 {shape_idx})")
                    return {"ocr_text": ""}
                else:
                    raise

            # 图片预处理用于 OCR 优化
            # 1. 灰度化
            image_data = image_data.convert('L')

            # 2. 对比度增强
            enhancer = ImageEnhance.Contrast(image_data)
            image_data = enhancer.enhance(2)

            # 3. Otsu 二值化
            np_image = np.array(image_data)
            _, np_image = cv2.threshold(
                np_image, 0, 255,
                cv2.THRESH_BINARY + cv2.THRESH_OTSU
            )
            image_data = Image.fromarray(np_image)

            # 4. OCR 识别
            ocr_text = pytesseract.image_to_string(
                image_data,
                lang='chi_sim+eng'  # 中英文混合识别
            )

            return {"ocr_text": ocr_text.strip()}

        except Exception as e:
            log_debug(f"⚠️ 图片 OCR 出错 (幻灯片 {slide_idx}, 形状 {shape_idx}): {str(e)}")
            return {"ocr_text": ""}

    def _parse_table(self, shape) -> Optional[dict]:
        """
        提取表格数据

        Args:
            shape: 表格形状对象

        Returns:
            dict: {"content": [[行1], [行2], ...]} 或 None
        """
        try:
            table_content = []
            for row in shape.table.rows:
                row_content = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_content.append(cell_text)
                table_content.append(row_content)

            return {"content": table_content}

        except Exception as e:
            log_debug(f"⚠️ 表格解析出错: {str(e)}")
            return None

    def _parse_group(self, shape, slide_idx: int, shape_idx, slide_content: dict):
        """
        递归处理分组形状

        Args:
            shape: 分组形状对象
            slide_idx: 幻灯片索引
            shape_idx: 形状索引
            slide_content: 存储解析结果的字典
        """
        for s_idx, s in enumerate(shape.shapes):
            self._parse_shape(s, slide_idx, f"{shape_idx}_{s_idx}", slide_content)

    def _table_to_text(self, table_content: List[List[str]]) -> str:
        """
        将表格数据转换为文本格式

        Args:
            table_content: 二维数组表格数据

        Returns:
            str: 文本格式的表格
        """
        if not table_content:
            return ""

        lines = []
        for row in table_content:
            line = " | ".join(row)
            lines.append(line)

        return "\n".join(lines)


# 1.2.42: 工厂函数 - 兼容 langchain 加载器接口
def create_pptx_loader(file_path: str, enable_ocr: bool = False) -> GeneralPPTXLoader:
    """
    创建 PPTX 加载器实例

    Args:
        file_path: PPT 文件路径
        enable_ocr: 是否启用图片 OCR

    Returns:
        GeneralPPTXLoader: 加载器实例
    """
    return GeneralPPTXLoader(file_path, enable_ocr=enable_ocr)
